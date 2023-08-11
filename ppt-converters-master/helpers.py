import base64
import glob
import re
import json
import uuid

from pptx import Presentation
RECOGNIZED_EXTENSIONS = ['png','jpg','bmp','jpeg','gif']

def pptToJson(INPUT_FILE):

    # Load presentation and create document file
    prs = Presentation(INPUT_FILE)
    file_name = uuid.uuid4()
    result={}

    # Image counter
    counter = 0
    total_duration = 0

    isFirstSlide = True
    total_slides = len(prs.slides)
    slide_number=0
    result["total_slides"]=total_slides
    result["stages_data"]={}

    i = 0
    for slide in prs.slides:
        paragraphs=[]
        slide_images=[]
        title =""
        if slide.shapes.title:
            title = slide.shapes.title.text
        result["stages_data"][str(slide_number)]={}
        result["stages_data"][str(slide_number)]["title"]=title

        result["stages_data"][str(slide_number)]["time"]="00:00:00"

        index = 0
        for shape in slide.shapes:


            if hasattr(shape, "animation_settings"):
                animation_settings = shape.animation_settings

                if hasattr(animation_settings, "advance"):
                    advance = animation_settings.advance
                    if advance.after == "0:00:00":
                        result["stages_data"][str(slide_number)]["time"]=advance.seconds
                        total_duration += advance.seconds
            if hasattr(shape, "transition"):
                transition = shape.transition
            if not shape.has_text_frame:
                # If no image then skip
                try:
                    shape.image
                except:
                    continue

                # Skip unrecognized formats
                if shape.image.ext not in RECOGNIZED_EXTENSIONS:
                    continue


                base64_image = base64.b64encode(shape.image.blob).decode('utf-8')
                slide_images.append(base64_image)
            if shape.has_text_frame:
                para_res={}
                for paragraph in shape.text_frame.paragraphs:
                    para_text=re.sub(u'[^\u0020-\uD7FF\u0009\u000A\u000D\uE000-\uFFFD\U00010000-\U0010FFFF]+', '', paragraph.text)

                    if paragraph.text!=title and para_text!="":
                        if index in para_res.keys():
                            print(index,paragraph.text)

                            para_res[index]=para_res[index]+"\n"+re.sub(u'[^\u0020-\uD7FF\u0009\u000A\u000D\uE000-\uFFFD\U00010000-\U0010FFFF]+', '', paragraph.text)
                        else:
                            print(index, paragraph.text)

                            para_res[index] = re.sub(
                                u'[^\u0020-\uD7FF\u0009\u000A\u000D\uE000-\uFFFD\U00010000-\U0010FFFF]+', '',
                                paragraph.text)

                paragraphs.append(para_res)

            index = index + 1
        result["stages_data"][str(slide_number)]["paragraphs"]=paragraphs
        result["stages_data"][str(slide_number)]["images"] = slide_images
        slide_number = slide_number + 1


    # prs.save("home/ifrahrao/mysite/static/" + str(file_name))
    # result["file_path"]="home/ifrahrao/mysite/static/" + str(file_name)
    # with open('home/ifrahrao/mysite/static/data.json', 'w', encoding='utf-8') as f:
    #     json.dump(result, f, ensure_ascii=False, indent=4)
    with open('output/data.json', 'w', encoding='utf-8') as f:
        json.dump(result, f, ensure_ascii=False, indent=4)
    return result

def change_content(presentation_path,slide_number,new_content,section_index):
    presentation = Presentation(presentation_path)
    try:

        # Get the specified slide and shape
        slide = presentation.slides[slide_number-1]
        shape = slide.shapes[section_index]
        if shape.has_text_frame:
            text_frame = shape.text_frame
            text_frame.clear()  # Clear existing content
            p = text_frame.add_paragraph()
            p.text = new_content


        # Save the modified presentation
        presentation.save("output/"+presentation_path.split("/")[-1])
        return "Content runtime updated successfully."
    except Exception as e:
        print(e)
        return e

def transform(a = 2):
    if a == 1:
        return a +- 2
    return a

total = 1

for x in [3,5,1]:
    total = total + transform(x)

print(total)

result=pptToJson("input/Review2.pptx")
change_content("input/Review2.pptx",1,"Souradeep Nanda 14BCE1037\nIfrah Nadeem14BCE1040\nPraveen Kumar 14BCE1117\nBarnik Bannerjee 14BCE1233",1)